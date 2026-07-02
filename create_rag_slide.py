"""
Génère une slide PPT du Pipeline RAG — éditable dans PowerPoint
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Couleurs
BLEU = RGBColor(0x15, 0x65, 0xC0)
VERT = RGBColor(0x2E, 0x7D, 0x32)
NOIR = RGBColor(0x21, 0x21, 0x21)
GRIS = RGBColor(0x61, 0x61, 0x61)
FOND_BEIGE = RGBColor(0xF5, 0xF2, 0xEB)

# Fond beige comme la capture
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = FOND_BEIGE

# ═══════════════════════════════════════
# TITRE
# ═══════════════════════════════════════
txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Pipeline RAG"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = NOIR
p.alignment = PP_ALIGN.LEFT

# ═══════════════════════════════════════
# 4 ÉTAPES (boîtes éditables)
# ═══════════════════════════════════════

etapes = [
    ("1. Ingestion & Fragmentation",
     "Extraction du texte (PDF, DOCX, TXT)\net découpage en blocs de 1000 caractères."),
    ("2. Vectorisation (Embeddings)",
     "Transformation en vecteurs via\nsentence-transformers et indexation dans ChromaDB."),
    ("3. Recherche Sémantique",
     "Calcul de similarité cosinus pour extraire\nles 5 passages les plus pertinents."),
    ("4. Génération de Réponse",
     "Le LLM (GPT-4o-mini) formule sa réponse\nuniquement à partir du contexte extrait."),
]

y_start = 3.0
espacement = 3.8

for i, (titre, desc) in enumerate(etapes):
    y = y_start + i * espacement

    # Boîte rectangle arrondi
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(2), Cm(y), Cm(12), Cm(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE5, 0xDE)
    shape.line.color.rgb = RGBColor(0xBD, 0xB9, 0xB0)
    shape.line.width = Pt(1)

    # Titre de l'étape
    txBox = slide.shapes.add_textbox(Cm(15), Cm(y), Cm(16), Cm(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titre
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NOIR

    # Description
    txBox2 = slide.shapes.add_textbox(Cm(15), Cm(y + 1.2), Cm(16), Cm(2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRIS

    # Flèche entre les étapes (sauf la dernière)
    if i < 3:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Cm(7.5), Cm(y + 3), Cm(1), Cm(0.7)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x9E, 0x9E, 0x9E)
        arrow.line.fill.background()

# ═══════════════════════════════════════
# NOTE EN BAS
# ═══════════════════════════════════════
txBox = slide.shapes.add_textbox(Cm(2), Cm(18), Cm(30), Cm(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Le LLM répond uniquement à partir du document réel — pas d'hallucination"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = VERT
p.alignment = PP_ALIGN.CENTER

prs.save('rfp-ai-platform/slide_rag_pipeline.pptx')
print("OK - slide_rag_pipeline.pptx")
